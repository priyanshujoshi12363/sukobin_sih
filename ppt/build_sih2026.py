"""
Builds the SIH 2026 idea-submission deck for Sukobin.

Follows the official six-slide structure: title, idea/solution, technical
approach, feasibility and viability, impact and benefits, research and
references.

Two rules this file follows:
  * Every claim is restricted to what is actually built and verified. The
    honest-status section of SUKOBIN_ALGORITHM.md is the source of truth.
  * Cards size themselves from their wrapped text. PowerPoint will not
    auto-grow a fixed rectangle, so text is wrapped here and the box height
    is derived from the resulting line count. Hard-coding heights is what
    made the first draft spill text past its borders.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import textwrap

# ── brand ────────────────────────────────────────────────────────────────────
FOREST = RGBColor(0x1A, 0x3D, 0x2B)
FOREST_DEEP = RGBColor(0x10, 0x26, 0x1B)
GREEN = RGBColor(0x0C, 0x83, 0x1F)
GREEN_MID = RGBColor(0x2D, 0x6A, 0x4F)
MINT = RGBColor(0xD8, 0xF3, 0xDC)
MINT_DIM = RGBColor(0x9E, 0xC9, 0xAB)
SAGE = RGBColor(0x7D, 0xAA, 0x90)
CREAM = RGBColor(0xF9, 0xF8, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1F, 0x29, 0x37)
GREY = RGBColor(0x6B, 0x72, 0x80)
GREY_LIGHT = RGBColor(0xE5, 0xE7, 0xEB)
AMBER = RGBColor(0xF4, 0xA2, 0x61)
RED = RGBColor(0xE6, 0x39, 0x46)

W = Inches(13.333)
H = Inches(7.5)
FONT = "Segoe UI"


def new_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.15):
    """runs: (string, size_pt, bold, color) with an optional 5th space_after."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    for i, r in enumerate(runs):
        content, size, bold, color = r[0], r[1], r[2], r[3]
        after = r[4] if len(r) > 4 else 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(after)
        run = p.add_run()
        run.text = content
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
    return box


# ── text metrics ─────────────────────────────────────────────────────────────
# Segoe UI averages ~0.48 em per character at these sizes. Wrapping here rather
# than letting PowerPoint do it means the line count is known before the box is
# drawn, so the box can be sized to fit.

def fit(lines, width_in, size, pad_in=0.44):
    # 0.54 em is deliberately pessimistic. If the estimate is even slightly
    # optimistic PowerPoint re-wraps the line and the card overflows.
    usable_pt = (width_in - pad_in) * 72.0
    chars = max(14, int(usable_pt / (0.54 * size)))
    out = []
    for ln in lines:
        if not ln.strip():
            out.append("")
            continue
        hanging = ln.startswith("   ")
        wrapped = textwrap.wrap(ln.strip(), width=chars if not hanging else chars - 3)
        if not wrapped:
            out.append("")
            continue
        out.append(("   " if hanging else "") + wrapped[0])
        out.extend("   " + w for w in wrapped[1:])
    return out


def line_h(size):
    # A rendered line is font_size * intrinsic_leading(1.2) * line_spacing,
    # plus the paragraph space_after. Omitting the 1.2 was why cards still
    # ran a line short.
    return (size * 1.2 * 1.24 + 3.0) / 72.0


def card_h(lines, width_in, size, title_size=13):
    n = len(fit(lines, width_in, size))
    return Inches(0.15 + title_size / 72.0 + 0.17 + n * line_h(size) + 0.16)


def card(slide, x, y, w, title, lines, accent=GREEN, title_size=13,
         body_size=10.5, fill=WHITE, h=None):
    width_in = w / 914400
    shown = fit(lines, width_in, body_size)
    if h is None:
        h = card_h(lines, width_in, body_size, title_size)

    rect(slide, x, y, w, h, fill=fill, line=GREY_LIGHT, line_w=0.75)
    rect(slide, x, y, Pt(4), h, fill=accent)

    text(slide, x + Inches(0.22), y + Inches(0.15), w - Inches(0.42), Inches(0.3),
         [(title, title_size, True, FOREST)])

    body_y = y + Inches(0.15 + title_size / 72.0 + 0.17)
    runs = [(ln, body_size, False, GREY, 3.0) for ln in shown]
    text(slide, x + Inches(0.22), body_y, w - Inches(0.42),
         Inches(len(shown) * line_h(body_size) + 0.2), runs, line_spacing=1.24)
    return h


def header(slide, number, title, subtitle=None):
    rect(slide, 0, 0, W, Inches(1.05), fill=FOREST)
    rect(slide, 0, Inches(1.05), W, Pt(3), fill=GREEN)
    text(slide, Inches(0.55), Inches(0.2), Inches(0.6), Inches(0.6),
         [(number, 26, True, SAGE)])
    text(slide, Inches(1.2), Inches(0.17), Inches(9.3), Inches(0.42),
         [(title, 25, True, WHITE)])
    if subtitle:
        text(slide, Inches(1.23), Inches(0.63), Inches(9.6), Inches(0.3),
             [(subtitle, 11, False, MINT_DIM)])
    text(slide, Inches(11.0), Inches(0.36), Inches(1.85), Inches(0.35),
         [("SUKOBIN", 12, True, MINT)], align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
def slide_title(prs):
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=FOREST)
    rect(s, Inches(8.75), 0, Inches(4.6), H, fill=FOREST_DEEP)

    text(s, Inches(0.85), Inches(0.85), Inches(7.5), Inches(0.4),
         [("SMART INDIA HACKATHON 2026", 13, True, SAGE)])

    text(s, Inches(0.85), Inches(1.45), Inches(7.6), Inches(1.1),
         [("SUKOBIN", 62, True, WHITE)], line_spacing=1.0)

    text(s, Inches(0.85), Inches(2.72), Inches(7.5), Inches(1.0),
         [("AI-Enabled Logistics Accessibility Intelligence", 23, True, MINT),
          ("Platform for the North Eastern Region", 23, True, MINT)],
         line_spacing=1.2)

    rect(s, Inches(0.85), Inches(4.05), Inches(0.7), Pt(3), fill=AMBER)

    text(s, Inches(0.85), Inches(4.4), Inches(7.3), Inches(1.5),
         [("The vehicles already making the journey carry the cargo —", 14, False, MINT_DIM),
          ("and because they never stop moving, they are also the", 14, False, MINT_DIM),
          ("sensor network that measures whether the road is passable.", 14, True, WHITE)],
         line_spacing=1.45)

    # proof strip fills what was dead space
    rect(s, Inches(0.85), Inches(6.25), Inches(7.3), Pt(1), fill=GREEN_MID)
    proof = [("42", "segments"), ("3,567 km", "network"),
             ("82", "districts"), ("4 + 1", "apps + dashboard")]
    for i, (v, l) in enumerate(proof):
        x = Inches(0.85) + i * Inches(1.85)
        text(s, x, Inches(6.5), Inches(1.7), Inches(0.35),
             [(v, 19, True, MINT)])
        text(s, x, Inches(6.83), Inches(1.7), Inches(0.28),
             [(l, 9.5, False, SAGE)])

    meta = [
        ("PROBLEM STATEMENT", ["SIH2026 — MDoNER"]),
        ("ORGANISATION", ["Ministry of Development of the", "North Eastern Region"]),
        ("THEME", ["Transportation & Logistics"]),
        ("CATEGORY", ["Software"]),
        ("TEAM NAME", ["< your team name >"]),
        ("TEAM ID", ["< your team id >"]),
    ]
    y = Inches(1.35)
    for label, values in meta:
        text(s, Inches(9.3), y, Inches(3.6), Inches(0.22),
             [(label, 8.5, True, SAGE)])
        y += Inches(0.26)
        for v in values:
            text(s, Inches(9.3), y, Inches(3.6), Inches(0.3),
                 [(v, 11.5, True, WHITE)])
            y += Inches(0.27)
        y += Inches(0.26)


# ─────────────────────────────────────────────────────────────────────────────
# 2 — Proposed solution
# ─────────────────────────────────────────────────────────────────────────────
def slide_idea(prs):
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=CREAM)
    header(s, "01", "Proposed Solution",
           "One platform that moves essential goods and measures the roads at the same time")

    top = Inches(1.38)

    h1 = card(s, Inches(0.5), top, Inches(3.95), "The problem, precisely",
              [
                  "NER districts lose road access to landslides, floods and snow.",
                  "Medicines and food are delayed, and nobody has a live picture of which roads are open.",
                  "A dedicated courier fleet is uneconomical here, so one was never built.",
              ], accent=RED, body_size=10)

    h2 = card(s, Inches(4.67), top, Inches(3.95), "The idea",
              [
                  "There is no fleet. Use the vehicles already travelling.",
                  "A tourist, commuter or taxi driver enters their number plate and route.",
                  "The app shows only the parcels riding that exact road, forward-facing, within capacity.",
              ], accent=GREEN, body_size=10)

    h3 = card(s, Inches(8.84), top, Inches(4.0), "The insight that answers the PS",
              [
                  "Those carriers stream GPS, so they do not only move cargo. They measure the road.",
                  "When vehicles that normally clear a ghat at 30 km/h all drop to 4, the system knows it is failing.",
              ], accent=AMBER, body_size=10)

    # the differentiator
    band_y = top + max(h1, h2, h3) + Inches(0.18)
    rect(s, Inches(0.5), band_y, Inches(12.34), Inches(1.18), fill=FOREST)
    text(s, Inches(0.8), band_y + Inches(0.15), Inches(11.7), Inches(0.26),
         [("WHAT IS NEW HERE", 9.5, True, SAGE)])
    text(s, Inches(0.8), band_y + Inches(0.45), Inches(11.7), Inches(0.65),
         [("Every other platform learns a road is blocked when somebody reports it.", 13, False, WHITE),
          ("Ours notices when six vehicles slow to a crawl — and it notices because those same vehicles are already carrying the medicines.", 13, True, MINT)],
         line_spacing=1.3)

    # clause coverage
    text(s, Inches(0.5), band_y + Inches(1.4), Inches(12), Inches(0.28),
         [("HOW IT ADDRESSES THE PROBLEM STATEMENT", 9.5, True, GREY)])

    pairs = [
        ("a  Accessibility monitoring", "Segment-level live status from carrier GPS and officer reports"),
        ("b  Disruption prediction", "Risk per segment from rainfall, terrain and past incidents"),
        ("c  Alternate routes and delay", "Alternates evaluated, blocked corridors refused, delay per segment"),
        ("d  GPS tracking of essentials", "Every consignment rides in a verified, position-streaming vehicle"),
        ("e  Automated alerts", "Blocked roads, risky corridors and delayed consignments pushed out"),
        ("f  Field reporting", "Geo-tagged photo reports; AI reads free text in local languages"),
        ("g  Central dashboards", "District connectivity, bottlenecks, emergency and supply views"),
        ("h  Multilingual and offline", "Reports queue offline and sync later; alerts in regional languages"),
    ]
    y0 = band_y + Inches(1.74)
    for i, (k, v) in enumerate(pairs):
        col, row = i % 2, i // 2
        x = Inches(0.5) + col * Inches(6.35)
        y = y0 + row * Inches(0.5)
        text(s, x, y, Inches(2.25), Inches(0.28), [(k, 10, True, FOREST)])
        wrapped = fit([v], 4.15, 9, pad_in=0.0)
        text(s, x + Inches(2.3), y, Inches(3.95), Inches(0.44),
             [(w, 9, False, GREY) for w in wrapped], line_spacing=1.2)


# ─────────────────────────────────────────────────────────────────────────────
# 3 — Technical approach
# ─────────────────────────────────────────────────────────────────────────────
def slide_technical(prs):
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=CREAM)
    header(s, "02", "Technical Approach",
           "Two independent evidence sources decide whether a road is passable — neither can act alone")

    # diagram panel
    rect(s, Inches(0.5), Inches(1.38), Inches(7.65), Inches(4.42),
         fill=WHITE, line=GREY_LIGHT, line_w=0.75)
    text(s, Inches(0.75), Inches(1.55), Inches(7.1), Inches(0.28),
         [("THE SENSING LOOP", 9.5, True, GREY)])

    for x, title, sub in [
        (Inches(0.85), "Carrier vehicles", "GPS every ~15 s while on a trip"),
        (Inches(4.62), "Field officers", "Geo-tagged photo + description"),
    ]:
        rect(s, x, Inches(1.95), Inches(3.15), Inches(0.7), fill=MINT, line=GREEN, line_w=1)
        text(s, x + Inches(0.15), Inches(2.06), Inches(2.85), Inches(0.5),
             [(title, 11.5, True, FOREST), (sub, 8.5, False, GREY)], line_spacing=1.25)

    for x, title, l1, l2 in [
        (Inches(0.85), "Map-match to segment", "Rolling median speed against", "the road's own baseline"),
        (Inches(4.62), "AI reads the report", "Local-language text becomes type,", "severity, blocks traffic, clearance"),
    ]:
        rect(s, x, Inches(2.95), Inches(3.15), Inches(0.9), fill=WHITE, line=GREY_LIGHT, line_w=0.75)
        text(s, x + Inches(0.15), Inches(3.06), Inches(2.85), Inches(0.7),
             [(title, 10.5, True, INK), (l1, 8.5, False, GREY), (l2, 8.5, False, GREY)],
             line_spacing=1.25)

    for cx in (Inches(2.42), Inches(6.19)):
        rect(s, cx, Inches(2.69), Pt(2.5), Inches(0.2), fill=GREEN)
        rect(s, cx, Inches(3.89), Pt(2.5), Inches(0.18), fill=GREEN)

    rect(s, Inches(0.85), Inches(4.09), Inches(6.92), Inches(0.72), fill=FOREST)
    text(s, Inches(1.05), Inches(4.2), Inches(6.5), Inches(0.5),
         [("Status resolver — confidence-weighted vote", 12, True, WHITE),
          ("Field report and probe evidence ranked, corroborated and capped", 8.5, False, MINT_DIM)],
         line_spacing=1.25)

    rect(s, Inches(4.28), Inches(4.85), Pt(2.5), Inches(0.18), fill=GREEN)

    for i, (t1, t2) in enumerate([
        ("Route matching", "blocked corridors refused"),
        ("Dashboard", "district colours, bottlenecks"),
        ("Alerts", "carriers, officers, citizens"),
    ]):
        x = Inches(0.85) + i * Inches(2.35)
        rect(s, x, Inches(5.07), Inches(2.15), Inches(0.55), fill=MINT)
        text(s, x + Inches(0.13), Inches(5.15), Inches(1.95), Inches(0.42),
             [(t1, 10, True, FOREST), (t2, 8, False, GREY)], line_spacing=1.2)

    # guards
    guards_h = card(s, Inches(8.4), Inches(1.38), Inches(4.44), "Guards that keep it honest",
                    [
                        "One vehicle cannot close a road: two distinct vehicles and four samples minimum.",
                        "GPS fixes worse than 120 m accuracy are discarded.",
                        "A weather forecast predicts risk; it never asserts that a road is currently shut.",
                        "An unverified report is capped at RESTRICTED until a human verifies it or probe data agrees.",
                    ], accent=AMBER, body_size=9.5)

    card(s, Inches(8.4), Inches(1.38) + guards_h + Inches(0.16), Inches(4.44), "Stack",
         [
             "Android — Kotlin and XML, four native apps on one shared core",
             "Backend — Node.js, Express, MongoDB with 2dsphere geo queries",
             "Routing — OSRM road polylines, alternates, corridor matching",
             "Weather — Open-Meteo hourly rainfall, snowfall, temperature",
             "AI — report classification with a deterministic keyword fallback",
             "Dashboard — React and MapLibre GIS",
             "Push — Firebase Cloud Messaging",
         ], accent=GREEN, body_size=9.5)

    # matching rule
    rect(s, Inches(0.5), Inches(5.95), Inches(7.65), Inches(1.05),
         fill=WHITE, line=GREY_LIGHT, line_w=0.75)
    text(s, Inches(0.75), Inches(6.08), Inches(7.1), Inches(0.26),
         [("THE MATCHING RULE", 9.5, True, GREY)])
    text(s, Inches(0.75), Inches(6.36), Inches(7.15), Inches(0.55),
         [("Pickup and drop must both lie inside the corridor of the driver's real", 9.5, False, INK),
          ("road polyline, forward-facing, with no blocked segment on the way.", 9.5, False, INK)],
         line_spacing=1.25)


# ─────────────────────────────────────────────────────────────────────────────
# 4 — Feasibility and viability
# ─────────────────────────────────────────────────────────────────────────────
def slide_feasibility(prs):
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=CREAM)
    header(s, "03", "Feasibility & Viability",
           "Built and running on live data, with the hard objections answered")

    rect(s, Inches(0.5), Inches(1.38), Inches(12.34), Inches(1.4), fill=FOREST)
    text(s, Inches(0.8), Inches(1.53), Inches(11.7), Inches(0.26),
         [("ALREADY BUILT AND RUNNING ON LIVE DATA — NOT A MOCK-UP", 9.5, True, SAGE)])

    for i, (v, l) in enumerate([
        ("42", "road segments modelled"),
        ("3,567 km", "on real road geometry"),
        ("82", "NER districts covered"),
        ("12", "supply chokepoints"),
        ("4 + 1", "native apps + dashboard"),
    ]):
        x = Inches(0.8) + i * Inches(2.4)
        text(s, x, Inches(1.9), Inches(2.3), Inches(0.42),
             [(v, 21, True, WHITE)], align=PP_ALIGN.CENTER)
        text(s, x, Inches(2.36), Inches(2.3), Inches(0.28),
             [(l, 9, False, MINT_DIM)], align=PP_ALIGN.CENTER)

    ha = card(s, Inches(0.5), Inches(3.0), Inches(6.1), "Verified behaviour, on real inputs",
         [
             "Risk model flagged Dimapur–Kohima SEVERE from 161 mm of actual recorded rainfall, a lifeline route for Manipur.",
             "Probe sensing graded a corridor OPEN, SLOW, RESTRICTED then BLOCKED as vehicle speeds fell.",
             "AI classified 8 of 8 messy Hinglish field reports correctly, extracting clearance time and which vehicles can still pass.",
             "The planner refused Dimapur to Imphal while NH-2 was cut.",
         ], accent=GREEN, body_size=9.5)

    hb = card(s, Inches(6.74), Inches(3.0), Inches(6.1), "Challenges and how they are handled",
         [
             "Trust — would you give medicines to a tourist? Tiered carriers: an open tier for commercial parcels, a trusted tier for essential cargo.",
             "Thin traffic on remote spurs — relay handoff at hubs, so a parcel changes carrier instead of waiting.",
             "No signal in the hills — offline queue that syncs on return.",
         ], accent=AMBER, body_size=9.5)

    band = Inches(3.0) + max(ha, hb) + Inches(0.16)
    rect(s, Inches(0.5), band, Inches(12.34), Inches(1.62),
         fill=WHITE, line=GREY_LIGHT, line_w=0.75)
    text(s, Inches(0.78), band + Inches(0.15), Inches(11.7), Inches(0.26),
         [("WHY IT SUSTAINS ITSELF", 9.5, True, GREY)])

    for i, (t1, t2) in enumerate([
        ("Zero fleet cost", "No vehicles to buy and no drivers to employ. Capacity already exists on the road."),
        ("Marginal delivery cost", "The journey happens regardless; the only real cost is the driver's small detour."),
        ("Data costs nothing extra", "The accessibility feed is a by-product of deliveries, so monitoring gets cheaper as usage grows."),
        ("Government-ready", "Cloud infrastructure, weather and transport API integration, district-level dashboards."),
    ]):
        x = Inches(0.78) + i * Inches(3.05)
        text(s, x, band + Inches(0.47), Inches(2.85), Inches(0.26), [(t1, 11, True, FOREST)])
        wrapped = fit([t2], 2.85, 9, pad_in=0.0)
        text(s, x, band + Inches(0.77), Inches(2.85), Inches(0.8),
             [(w, 9, False, GREY) for w in wrapped], line_spacing=1.25)


# ─────────────────────────────────────────────────────────────────────────────
# 5 — Impact and benefits
# ─────────────────────────────────────────────────────────────────────────────
def slide_impact(prs):
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=CREAM)
    header(s, "04", "Impact & Benefits",
           "What changes for a district officer, a family in a remote block, and a person with a car")

    rect(s, Inches(0.5), Inches(1.38), Inches(12.34), Inches(1.72), fill=FOREST)
    text(s, Inches(0.8), Inches(1.53), Inches(11.7), Inches(0.26),
         [("WHAT THE PLATFORM DOES WHEN A ROAD CLOSES", 9.5, True, SAGE)])

    steps = [
        ("1", "Rain rises", ["Risk on the ghat", "turns severe"]),
        ("2", "Officer reports", ["Photo and text, filed", "offline at the site"]),
        ("3", "AI reads it", ["Landslide, critical,", "road fully blocked"]),
        ("4", "Verified", ["Senior officer", "confirms in minutes"]),
        ("5", "Road marked shut", ["Probe data agrees:", "no vehicle moving"]),
        ("6", "System reacts", ["Carriers rerouted,", "district flagged"]),
    ]
    for i, (n, t1, t2) in enumerate(steps):
        x = Inches(0.8) + i * Inches(2.02)
        text(s, x, Inches(1.9), Inches(0.3), Inches(0.3), [(n, 15, True, AMBER)])
        text(s, x + Inches(0.28), Inches(1.93), Inches(1.68), Inches(0.26),
             [(t1, 10.5, True, WHITE)])
        for j, line in enumerate(t2):
            text(s, x, Inches(2.32) + j * Inches(0.21), Inches(1.9), Inches(0.22),
                 [(line, 8.5, False, MINT_DIM)])
        if i < len(steps) - 1:
            rect(s, x + Inches(1.86), Inches(1.98), Pt(2), Inches(0.55), fill=GREEN_MID)

    benefits = [
        ("Social", RED, [
            "Medicines and food reach blocks that couriers do not serve.",
            "Emergency response gets a live map instead of phone calls.",
            "Alerts arrive in the language the recipient reads.",
        ]),
        ("Economic", GREEN, [
            "Travellers earn on a journey they already make.",
            "Hill shopkeepers reach customers beyond their valley.",
            "Districts avoid the cost of stock-outs and stranded goods.",
        ]),
        ("Environmental", GREEN_MID, [
            "No new fleet, so no additional vehicles on the road.",
            "Deliveries ride in vehicles already making the trip.",
            "Fewer empty return legs across mountain corridors.",
        ]),
        ("Governance", AMBER, [
            "One accessibility picture across all NER districts.",
            "Field reports become structured, auditable records.",
            "Disruption is predicted, not recorded after the fact.",
        ]),
    ]
    for i, (title, accent, lines) in enumerate(benefits):
        x = Inches(0.5) + i * Inches(3.13)
        card(s, x, Inches(3.3), Inches(2.98), title, lines,
             accent=accent, body_size=9, h=Inches(2.05))

    rect(s, Inches(0.5), Inches(5.55), Inches(12.34), Inches(1.1),
         fill=WHITE, line=GREEN, line_w=1.5)
    text(s, Inches(0.85), Inches(5.75), Inches(11.7), Inches(0.7),
         [("Sukobin turns every journey already being made into delivery capacity — and every carrier into a road sensor.", 15, True, FOREST),
          ("The network that carries the medicines is the same network that tells the government whether the road is open.", 12, False, GREY)],
         line_spacing=1.4)


# ─────────────────────────────────────────────────────────────────────────────
# 6 — Research and references
# ─────────────────────────────────────────────────────────────────────────────
def slide_research(prs):
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=CREAM)
    header(s, "05", "Research & References",
           "Sources, data providers and the corridors the system is modelled on")

    ch = Inches(2.28)

    card(s, Inches(0.5), Inches(1.38), Inches(6.1), "Data sources and APIs",
         [
             "OSRM — road network routing and alternate paths",
             "Open-Meteo — hourly rainfall, snowfall, temperature and elevation, no key required",
             "OpenStreetMap and Esri — base cartography",
             "Vahan — vehicle registration verification for carriers",
             "Firebase Cloud Messaging — multilingual push delivery",
             "MongoDB 2dsphere — geospatial corridor queries",
         ], accent=GREEN, body_size=9.5, h=ch)

    card(s, Inches(6.74), Inches(1.38), Inches(6.1), "NER corridors modelled",
         [
             "NH-10 Siliguri–Gangtok — Sikkim's lifeline, Teesta slides",
             "NH-2 Dimapur–Kohima–Imphal — Manipur's main supply line",
             "NH-6 Shillong–Jowai–Silchar — cuts off the Barak Valley",
             "NH-27 Siliguri–Guwahati — the Siliguri Corridor",
             "NH-13 Tezpur–Sela Pass–Tawang — snow closure Dec to Mar",
             "NH-306 Silchar–Aizawl — Mizoram's only all-weather road",
             "Plus NH-8, NH-715, NH-37, NH-127B, NH-29 — twelve in total",
         ], accent=AMBER, body_size=9.5, h=ch)

    ch2 = Inches(1.95)

    card(s, Inches(0.5), Inches(3.82), Inches(6.1), "Problem context",
         [
             "MDoNER — North East vision and infrastructure programmes",
             "NHIDCL and BRO — highway and border road status in NER",
             "NDMA — landslide and flood hazard guidance for the region",
             "IMD — rainfall thresholds behind slope-failure reasoning",
             "PMGSY — rural road connectivity for last-mile blocks",
         ], accent=FOREST, body_size=9.5, h=ch2)

    card(s, Inches(6.74), Inches(3.82), Inches(6.1), "Technical foundations",
         [
             "Probe-vehicle traffic estimation — inferring road condition from vehicle speed traces",
             "Rainfall-threshold landslide models — antecedent rainfall as the dominant trigger on steep slopes",
             "Crowdsourced logistics — capacity sharing on existing trips",
             "Idempotent offline sync — client-generated keys for safe replay",
         ], accent=GREY, body_size=9.5, h=ch2)

    rect(s, Inches(0.5), Inches(6.02), Inches(12.34), Inches(0.8), fill=FOREST)
    text(s, Inches(0.85), Inches(6.22), Inches(11.7), Inches(0.42),
         [("Working code, live data, and a demo that runs end to end — repository and builds available on request.", 12, True, MINT)])


def main():
    prs = new_deck()
    for fn in (slide_title, slide_idea, slide_technical,
               slide_feasibility, slide_impact, slide_research):
        fn(prs)

    out = os.path.normpath(
        os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "SIH2026_Sukobin.pptx")
    )
    prs.save(out)
    print("written:", out)


if __name__ == "__main__":
    main()
